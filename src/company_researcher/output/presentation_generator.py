"""
Presentation Generator (Phase 16.3).

PowerPoint/Slides generation:
- Executive presentations
- Investor pitches
- Board updates
- Research summaries

Supports python-pptx with markdown fallback.
"""

import os
from dataclasses import dataclass, field
from enum import Enum
from typing import Any, Dict, List, Optional

from ..utils import utc_now

# Try to import presentation library
try:
    from pptx import Presentation
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ============================================================================
# Data Models
# ============================================================================


class PresentationStyle(str, Enum):
    """Presentation style options."""

    EXECUTIVE = "executive"  # Brief, high-level
    DETAILED = "detailed"  # Comprehensive
    INVESTOR = "investor"  # Investment focused
    BOARD = "board"  # Board meeting format


class SlideType(str, Enum):
    """Types of slides."""

    TITLE = "title"
    SECTION = "section"
    CONTENT = "content"
    TWO_COLUMN = "two_column"
    METRICS = "metrics"
    CHART = "chart"
    BULLETS = "bullets"
    CONCLUSION = "conclusion"


@dataclass
class SlideContent:
    """Content for a single slide."""

    slide_type: SlideType
    title: str
    subtitle: str = ""
    body: str = ""
    bullets: List[str] = field(default_factory=list)
    metrics: Dict[str, Any] = field(default_factory=dict)
    notes: str = ""


@dataclass
class PresentationConfig:
    """Configuration for presentation generation."""

    title: str
    company_name: str
    style: PresentationStyle = PresentationStyle.EXECUTIVE
    max_slides: int = 20
    include_appendix: bool = False
    author: str = "Company Researcher"
    date: str = field(default_factory=lambda: utc_now().strftime("%Y-%m-%d"))


# ============================================================================
# Presentation Generator
# ============================================================================


class PresentationGenerator:
    """
    Generate PowerPoint presentations from research data.

    Usage:
        generator = PresentationGenerator()
        path = generator.generate(
            research_data=data,
            config=PresentationConfig(
                title="Tesla Analysis",
                company_name="Tesla",
                style=PresentationStyle.INVESTOR
            )
        )
    """

    # Color scheme
    COLORS = {
        "primary": "2C5282",
        "secondary": "2B6CB0",
        "accent": "38A169",
        "text": "2D3748",
        "light": "EDF2F7",
    }

    def __init__(self):
        """Initialize generator."""
        self._pptx_available = PPTX_AVAILABLE

    def generate(
        self,
        research_data: Dict[str, Any],
        config: PresentationConfig,
        output_path: Optional[str] = None,
    ) -> str:
        """
        Generate presentation.

        Args:
            research_data: Research data from agents
            config: Presentation configuration
            output_path: Output file path

        Returns:
            Path to generated presentation
        """
        if not output_path:
            safe_name = config.company_name.lower().replace(" ", "_")
            output_path = f"./reports/{safe_name}_presentation_{config.date}.pptx"

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

        if self._pptx_available:
            return self._generate_pptx(research_data, config, output_path)
        else:
            return self._generate_markdown_fallback(research_data, config, output_path)

    def _generate_pptx(
        self, data: Dict[str, Any], config: PresentationConfig, output_path: str
    ) -> str:
        """Generate PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Generate slides
        slides = self._create_slide_content(data, config)

        for slide_content in slides[: config.max_slides]:
            self._add_slide(prs, slide_content)

        # Save presentation
        prs.save(output_path)
        return output_path

    def _create_slide_content(
        self, data: Dict[str, Any], config: PresentationConfig
    ) -> List[SlideContent]:
        """Create slide content from research data."""
        slides = []
        agent_outputs = data.get("agent_outputs", data)

        # Title slide
        slides.append(
            SlideContent(
                slide_type=SlideType.TITLE,
                title=config.title,
                subtitle=f"{config.company_name} | {config.date}",
            )
        )

        # Executive summary
        slides.append(SlideContent(slide_type=SlideType.SECTION, title="Executive Summary"))

        key_findings = self._extract_key_findings(agent_outputs)
        if key_findings:
            slides.append(
                SlideContent(
                    slide_type=SlideType.BULLETS, title="Key Findings", bullets=key_findings[:6]
                )
            )

        # Key metrics slide
        metrics = self._extract_metrics(agent_outputs)
        if metrics:
            slides.append(
                SlideContent(slide_type=SlideType.METRICS, title="Key Metrics", metrics=metrics)
            )

        # Financial section
        if "financial" in agent_outputs:
            slides.append(SlideContent(slide_type=SlideType.SECTION, title="Financial Analysis"))
            slides.append(self._create_financial_slide(agent_outputs["financial"]))

        # Market section
        if "market" in agent_outputs:
            slides.append(SlideContent(slide_type=SlideType.SECTION, title="Market Analysis"))
            slides.append(self._create_market_slide(agent_outputs["market"]))

        # Competitive section
        if "competitor" in agent_outputs:
            slides.append(SlideContent(slide_type=SlideType.SECTION, title="Competitive Landscape"))
            slides.append(self._create_competitive_slide(agent_outputs["competitor"]))

        # Investment section (for investor style)
        if config.style == PresentationStyle.INVESTOR and "investment" in agent_outputs:
            slides.append(SlideContent(slide_type=SlideType.SECTION, title="Investment Analysis"))
            slides.append(self._create_investment_slide(agent_outputs["investment"]))

        # Recommendations
        recommendations = self._extract_recommendations(agent_outputs)
        if recommendations:
            slides.append(
                SlideContent(
                    slide_type=SlideType.BULLETS,
                    title="Recommendations",
                    bullets=recommendations[:6],
                )
            )

        # Conclusion slide
        slides.append(
            SlideContent(
                slide_type=SlideType.CONCLUSION,
                title="Conclusion",
                bullets=self._create_conclusion_points(agent_outputs),
            )
        )

        return slides

    def _add_slide(self, prs, content: SlideContent):
        """Add a slide to the presentation."""
        if content.slide_type == SlideType.TITLE:
            self._add_title_slide(prs, content)
        elif content.slide_type == SlideType.SECTION:
            self._add_section_slide(prs, content)
        elif content.slide_type == SlideType.BULLETS:
            self._add_bullet_slide(prs, content)
        elif content.slide_type == SlideType.METRICS:
            self._add_metrics_slide(prs, content)
        elif content.slide_type == SlideType.CONTENT:
            self._add_content_slide(prs, content)
        elif content.slide_type == SlideType.CONCLUSION:
            self._add_conclusion_slide(prs, content)
        else:
            self._add_content_slide(prs, content)

    def _add_title_slide(self, prs, content: SlideContent):
        """Add title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor.from_string(self.COLORS["primary"])
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        if content.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = content.subtitle
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RgbColor.from_string(self.COLORS["text"])
            subtitle_para.alignment = PP_ALIGN.CENTER

    def _add_section_slide(self, prs, content: SlideContent):
        """Add section divider slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Section title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor.from_string(self.COLORS["secondary"])
        title_para.alignment = PP_ALIGN.CENTER

    def _add_bullet_slide(self, prs, content: SlideContent):
        """Add bullet point slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor.from_string(self.COLORS["primary"])

        # Bullets
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True

        for i, bullet in enumerate(content.bullets):
            if i == 0:
                para = body_frame.paragraphs[0]
            else:
                para = body_frame.add_paragraph()

            para.text = f"• {bullet}"
            para.font.size = Pt(20)
            para.font.color.rgb = RgbColor.from_string(self.COLORS["text"])
            para.space_after = Pt(12)

    def _add_metrics_slide(self, prs, content: SlideContent):
        """Add metrics slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor.from_string(self.COLORS["primary"])

        # Metrics in grid
        metrics = list(content.metrics.items())[:6]
        cols = 3
        rows = (len(metrics) + cols - 1) // cols

        for i, (metric_name, metric_value) in enumerate(metrics):
            row = i // cols
            col = i % cols

            x = Inches(0.5 + col * 4.2)
            y = Inches(1.8 + row * 2)

            # Metric value
            value_box = slide.shapes.add_textbox(x, y, Inches(4), Inches(0.8))
            value_frame = value_box.text_frame
            value_para = value_frame.paragraphs[0]
            value_para.text = str(metric_value)
            value_para.font.size = Pt(36)
            value_para.font.bold = True
            value_para.font.color.rgb = RgbColor.from_string(self.COLORS["secondary"])
            value_para.alignment = PP_ALIGN.CENTER

            # Metric name
            name_box = slide.shapes.add_textbox(x, Inches(y.inches + 0.8), Inches(4), Inches(0.4))
            name_frame = name_box.text_frame
            name_para = name_frame.paragraphs[0]
            name_para.text = metric_name
            name_para.font.size = Pt(16)
            name_para.font.color.rgb = RgbColor.from_string(self.COLORS["text"])
            name_para.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs, content: SlideContent):
        """Add general content slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor.from_string(self.COLORS["primary"])

        # Body
        if content.body:
            body_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
            )
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            body_para = body_frame.paragraphs[0]
            body_para.text = content.body[:1000]
            body_para.font.size = Pt(18)
            body_para.font.color.rgb = RgbColor.from_string(self.COLORS["text"])

    def _add_conclusion_slide(self, prs, content: SlideContent):
        """Add conclusion slide."""
        self._add_bullet_slide(prs, content)

    def _create_financial_slide(self, data: Dict[str, Any]) -> SlideContent:
        """Create financial analysis slide content."""
        bullets = []

        if "revenue" in data:
            bullets.append(f"Revenue: {data['revenue']}")
        if "growth_rate" in data:
            bullets.append(f"Growth Rate: {data['growth_rate']}")
        if "gross_margin" in data:
            bullets.append(f"Gross Margin: {data['gross_margin']}")

        analysis = data.get("analysis", "")
        if analysis and len(bullets) < 4:
            # Extract key sentences
            sentences = analysis.split(". ")[:3]
            for s in sentences:
                if len(s) > 20 and len(s) < 150:
                    bullets.append(s.strip())

        return SlideContent(
            slide_type=SlideType.BULLETS, title="Financial Highlights", bullets=bullets[:6]
        )

    def _create_market_slide(self, data: Dict[str, Any]) -> SlideContent:
        """Create market analysis slide content."""
        bullets = []

        if "tam" in data:
            bullets.append(f"TAM: {data['tam']}")
        if "sam" in data:
            bullets.append(f"SAM: {data['sam']}")
        if "som" in data:
            bullets.append(f"SOM: {data['som']}")
        if "market_share" in data:
            bullets.append(f"Market Share: {data['market_share']}")

        return SlideContent(
            slide_type=SlideType.BULLETS, title="Market Opportunity", bullets=bullets[:6]
        )

    def _create_competitive_slide(self, data: Dict[str, Any]) -> SlideContent:
        """Create competitive analysis slide content."""
        bullets = []

        if "competitors_found" in data:
            bullets.append(f"Competitors Identified: {data['competitors_found']}")
        if "competitive_intensity" in data:
            bullets.append(f"Competitive Intensity: {data['competitive_intensity']}")
        if "threat_summary" in data:
            bullets.append(f"Threat Assessment: {data['threat_summary']}")

        return SlideContent(
            slide_type=SlideType.BULLETS, title="Competitive Position", bullets=bullets[:6]
        )

    def _create_investment_slide(self, data: Dict[str, Any]) -> SlideContent:
        """Create investment analysis slide content."""
        bullets = []

        if "investment_rating" in data:
            bullets.append(f"Rating: {data['investment_rating'].upper()}")
        if "overall_risk" in data:
            bullets.append(f"Risk Level: {data['overall_risk']}")
        if "moat_strength" in data:
            bullets.append(f"Competitive Moat: {data['moat_strength']}")
        if "growth_stage" in data:
            bullets.append(f"Growth Stage: {data['growth_stage']}")

        return SlideContent(
            slide_type=SlideType.BULLETS, title="Investment Assessment", bullets=bullets[:6]
        )

    def _extract_key_findings(self, data: Dict[str, Any]) -> List[str]:
        """Extract key findings from all agent outputs."""
        findings = []

        for key, value in data.items():
            if isinstance(value, dict):
                if "key_findings" in value:
                    findings.extend(value["key_findings"][:2])
                elif "conclusions" in value:
                    findings.extend(value["conclusions"][:2])

        return findings[:6]

    def _extract_metrics(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Extract key metrics for display."""
        metrics = {}

        if "financial" in data:
            if "revenue" in data["financial"]:
                metrics["Revenue"] = data["financial"]["revenue"]

        if "market" in data:
            if "market_share" in data["market"]:
                metrics["Market Share"] = data["market"]["market"]

        if "brand" in data:
            if "brand_score" in data["brand"]:
                metrics["Brand Score"] = data["brand"]["brand_score"]

        if "investment" in data:
            if "investment_rating" in data["investment"]:
                metrics["Rating"] = data["investment"]["investment_rating"].upper()

        return metrics

    def _extract_recommendations(self, data: Dict[str, Any]) -> List[str]:
        """Extract recommendations from all agents."""
        recommendations = []

        for key, value in data.items():
            if isinstance(value, dict) and "recommendations" in value:
                recommendations.extend(value["recommendations"][:2])

        return recommendations[:6]

    def _create_conclusion_points(self, data: Dict[str, Any]) -> List[str]:
        """Create conclusion bullet points."""
        points = [
            "Comprehensive research analysis completed",
            "Multiple data sources analyzed",
            "Strategic insights generated",
        ]

        if "financial" in data:
            points.append("Financial health assessed")
        if "market" in data:
            points.append("Market opportunity quantified")
        if "competitor" in data:
            points.append("Competitive landscape mapped")

        return points[:6]

    def _generate_markdown_fallback(
        self, data: Dict[str, Any], config: PresentationConfig, output_path: str
    ) -> str:
        """Generate markdown slides as fallback."""
        md_path = output_path.replace(".pptx", "_slides.md")

        with open(md_path, "w") as f:
            f.write(f"# {config.title}\n\n")
            f.write(f"**{config.company_name}** | {config.date}\n\n")
            f.write("---\n\n")

            slides = self._create_slide_content(data, config)

            for slide in slides:
                if slide.slide_type == SlideType.TITLE:
                    f.write(f"# {slide.title}\n\n")
                    if slide.subtitle:
                        f.write(f"*{slide.subtitle}*\n\n")
                elif slide.slide_type == SlideType.SECTION:
                    f.write(f"## {slide.title}\n\n")
                elif slide.slide_type == SlideType.BULLETS:
                    f.write(f"### {slide.title}\n\n")
                    for bullet in slide.bullets:
                        f.write(f"- {bullet}\n")
                    f.write("\n")
                elif slide.slide_type == SlideType.METRICS:
                    f.write(f"### {slide.title}\n\n")
                    for metric, value in slide.metrics.items():
                        f.write(f"- **{metric}**: {value}\n")
                    f.write("\n")

                f.write("---\n\n")

        return md_path


# ============================================================================
# Factory Function
# ============================================================================


def generate_presentation(
    research_data: Dict[str, Any],
    company_name: str,
    title: Optional[str] = None,
    style: str = "executive",
    output_path: Optional[str] = None,
) -> str:
    """
    Generate presentation from research data.

    Args:
        research_data: Research data from agents
        company_name: Company name
        title: Presentation title
        style: Presentation style
        output_path: Output file path

    Returns:
        Path to generated presentation
    """
    config = PresentationConfig(
        title=title or f"{company_name} Research Analysis",
        company_name=company_name,
        style=PresentationStyle(style),
    )

    generator = PresentationGenerator()
    return generator.generate(research_data, config, output_path)
